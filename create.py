from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import math

img_path = 'sample.png'

def replace_txt(replacements, shapes):
    for shape in shapes:
        for match, replacement in replacements.items():
            #find tags inside text frame
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(match), str(replacement))
                            run.text = new_text
            #find tags inside table frame
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if match in cell.text:
                            new_text = cell.text.replace(match, replacement)
                            cell.text = new_text

# def _get_blank_slide_layout(pres):
#          layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
#          min_items = min(layout_items_count)
#          blank_layout_id = layout_items_count.index(min_items)
#          return pres.slide_layouts[blank_layout_id]

# def copy_slide(pres,pres1,index):
#         source = pres.slides[index]

#         blank_slide_layout = _get_blank_slide_layout(pres)
#         dest = pres1.slides.add_slide(blank_slide_layout)

#         for shp in source.shapes:
#             el = shp.element
#             newel = copy.deepcopy(el)
#             dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

#             for key, value in six.iteritems(source.rels):
#                         # Make sure we don't copy a notesSlide relation as that won't exist
#                     if not "notesSlide" in value.reltype:
#                             dest.rels.add_relationship(value.reltype, value._target, value.rId)

#             return dest



def replace_images(replacements, slide, index):
    if str(index) in replacements:
        objectKey = str(index)
        replacedData = replacements[objectKey]
        print(replacedData)

        for data in replacedData:
            print(data)
            slide.shapes.add_picture(data["path"], data["size"]["left"], data["size"]["top"], data["size"]["width"] ,data["size"]["height"] )

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')


def tables_headers(table, headers):
    i = 0
    for header in headers:
        table.cell(0, i).text = header
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(200, 253, 251)
        _set_cell_border(cell,"949595", '12000')
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(5)
        para.font.name = 'Comic Sans MS'
        para.font.color.rgb = RGBColor(19, 170, 246)
        table.columns[i].width = Inches(0.6)
        i += 1

def tables_rows(table, rowData, start, end,totalRows):
    j = start
    cell_start = 1

    entCount = end
    if end > totalRows:
        entCount = totalRows

    while j < entCount + 1:
        element = rowData[j - 1]
        k = 0
        if element:
            for key, value in element.items():
                table.cell(cell_start, k).text = value
                cell = table.cell(cell_start, k)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(228, 228, 228)
                _set_cell_border(cell,"949595", '12000')
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(6)
                para.font.name = 'Comic Sans MS'
                k += 1
        j += 1
        cell_start += 1


# def add_extra_data_to_next_slide():

#  remove unwanted templates 
def delete_slides(presentation, index):
    xml_slides = presentation.slides._sldIdLst  
    slides = list(xml_slides)
    print(index)
    try:
        slides[index]
        xml_slides.remove(slides[index])
    except ValueError:
        print("error") 
        

def replace_tables(replacements, slide , presentation, slideIndex):
    top, width, height  = Inches(0.6), Inches(2), Inches(0.04)
    rows = replacements["rowCount"]
    cols = replacements["columCount"]
    headers = replacements["headers"]
    rowData = replacements["rows"]
    extraSlideIndex = replacements["extraSlideIndex"]
    totalRows = len(rowData)
    totaltableCount = math.ceil(totalRows / rows)

    slideCount = math.ceil(totaltableCount / 4) 

    extraSlideCount = slideCount - 1

    if(extraSlideCount > 0):
        i = 0
        while i < extraSlideCount:
            extraSlideIndex.pop(0)
            i += 1

    loop = 0
    for preInedx in extraSlideIndex:
        deleteIndex = preInedx
        if(loop > 0):
            deleteIndex -= 1
        delete_slides(presentation, deleteIndex)
        loop += 1

    s = 0
    end = 0

    while s < slideCount:
        j = 0
        left = 1
        if totaltableCount > 4 :
            tableCount = 4
        else:
            tableCount = totaltableCount

        slideRowStart = 0
        slideRowEnd = 0
        while j < tableCount:
            shape = presentation.slides[slideIndex + s].shapes.add_table(rows + 1, cols, Inches(left) , top, width, height)
            table = shape.table
            tables_headers(table, headers)

            slideRowStart = (rows * j ) + 1 + end
            slideRowEnd = slideRowStart + (rows - 1)

            tables_rows(table,rowData, slideRowStart, slideRowEnd ,totalRows )
            left += 2
            j += 1

        end = slideRowEnd
        totaltableCount -= 4
        s += 1


if __name__ == '__main__':

    prs = Presentation('input.pptx')
    dataBump = {
        "image_replaces": {
                        "1": [{"path": "1.png", "size":{"left":Inches(2),"top":Inches(2), "height":Inches(3), "width":Inches(8) }}]
                    },
        "text_replaces":  {
            '{{var1}}': 'LGIM Bridge',
            '{{var2}}': 'External asset workflow',
            '{{var3}}': 'Active Client page – External assets',
            '{{var4}}': 'LGIM/Client user',
            '{{var5}}': 'Active Client page – External assets',
            '{{var6}}': 'LGIM/Client user – when the user clicks on ‘Edit detail’ for a date',
            '{{var7}}': 'Active Client page – External assets',
            '{{var8}}': 'Information button text',
            '{{var10}}': "If the scheme holds assets which are not managed by LGIM, details can be entered here. This information will be used for tracking the funding level, and also for overall portfolio risk analysis and funding level projections.Enter the value of non-LGIM assets held at different dates. We will roll the asset value forward in an approximate manner between dates provided.By default any asset values provided will be treated as cash for both funding level tracking and risk analysis purposes. However, if you click on ‘Edit details’ for a given asset value then you will be able to provide a breakdown of the value by fund and asset class. If this information is provided then we will roll forward the asset value in line with index returns for the relevant asset class, and will include the external assets within overall portfolio risk analysis. Once you have entered the breakdown of the asset value at one date, this breakdown will be carried forward to subsequent dates by default, unless you choose to edit the detail within these later dates (for example, following a significant change in the portfolio). This enables you to quickly update the value of non-LGIM assets without having to re-enter the underlying detail at every date.",
            '{{var11}}': 'PV01 and IE01 data should be available from your LDI provider. Note that positive numbers should typically be entered.',
            '{{var12}}': 'PV01 - Please enter the increase in asset value that would expected in £ terms if interest rates fell 0.01%.',
            '{{var13}}':'IE01 - Please enter the increase in asset value that would expected in £ terms if price inflation expectations increased 0.01%.',
            '{{var14}}':'If you do not know this information then you can simply leave these boxes blank, and we will treat the LDI assets as cash.'
            },
        "table_replaces": {
            "cashFlows":{
                "columCount": 3,
                "rowCount": 30,
                "headers": ["cashflow year","cashflow fixed","cashflow real"],
                "extraSlideIndex": [ 6,7,8 ],
                "rows": [ 
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"2","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"3","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"4","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"6","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"7","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"8","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"9","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"10","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"11","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"12","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"13","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"14","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"15","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"16","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"17","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"18","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"19","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"20","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"21","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"22","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"23","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"24","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"25","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"26","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"27","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"28","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"29","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"30","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"31","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"32","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"33","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"34","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"35","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"36","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"37","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"38","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"39","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"40","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"41","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"42","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"43","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"44","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"45","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"51","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"52","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"53","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"54","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"55","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"56","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"57","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"58","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"59","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"510","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"151","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"512","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"153","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"154","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"155","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"61","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"62","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"63","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"64","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"65","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"66","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"67","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"68","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"69","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"170","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"171","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"172","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"173","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"174","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"175","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"17","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"82","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"83","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"84","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"85","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"86","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"87","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"88","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"89","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"190","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"191","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"192","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"193","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"194","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"195","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"17","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"82","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"83","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"84","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"85","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"86","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"87","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"88","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"89","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"190","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"191","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"192","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"193","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"194","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"195","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"2","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"3","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"4","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"6","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"7","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"8","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"9","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"10","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"11","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"12","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"13","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"14","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"15","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"16","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"17","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"18","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"19","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"20","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"21","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"22","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"23","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"24","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"25","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"26","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"27","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"28","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"29","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"30","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"31","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"32","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"33","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"34","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"35","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"36","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"37","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"38","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"39","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"40","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"41","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"42","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"43","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"44","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"45","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"51","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5112","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5113","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5114","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5115","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5116","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5117","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"58","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5911","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"51110","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"15111","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5112","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1513","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1514","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1515","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"611","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"612","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"6113","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"614","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"615","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"616","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"617","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"618","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"619","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1170","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1171","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1172","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1173","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1714","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1715","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1117","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"812","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"813","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"814","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"815","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"816","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"817","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"818","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"819","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1910","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1911","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1912","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1913","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1914","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1915","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"117","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"812","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"813","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"814","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"815","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"816","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"187","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"818","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"189","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1190","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1911","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1192","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1913","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1941","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1915","cashflow_real": "123123"
                    }
                 ]
            }
        } 
             
    }

    slides = [slide for slide in prs.slides]
    shapes = []
    for slide in slides:
        print(slide)
        if "image_replaces" in dataBump:
            replace_images(dataBump["image_replaces"], slide, slides.index(slide))

        if "text_replaces" in dataBump:
            for shape in slide.shapes:
                shapes.append(shape)

        if(slides.index(slide) == 5):
            replace_tables(dataBump["table_replaces"]["cashFlows"], slide, prs,slides.index(slide))

    if "text_replaces" in dataBump:
        replace_txt(dataBump["text_replaces"], shapes)

    prs.save('output.pptx')
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import math
import re
import pydash
import time
from helpers.utils import check_tag_exist, replace_tags, get_tag_content, remove_extra_slides, drow_toc
from helpers.texts import text_replace
from helpers.images import replace_images
from helpers.tables import replace_tables, update_table_text, drow_tables

from expressions.for_loop import looper
from expressions.if_condition import _if

if __name__ == '__main__':
    start_time = time.perf_counter ()
    prs = Presentation('demo.pptx')

    replacements = {
        "schemeName": "XYZ Pension Scheme",
        "title": "Q2 2021 Summary Report",
        "heading":"Investment performance to 30 June 2021",
        "heading_assets":"Investment performance to 30 June 2021",
        "assetAllocation": "Asset allocation at 30 June 2021",
        "assetChart": { "url" : "img1.png" , "size": {"left":1,"top":1, "height":3, "width":4.2}},
        "overallPerformance": "Overall performance",
        "OPData":{
            "styles" : {
                "rw_1": {
                    "font_size": 9,
                    "font_name": "Arial",
                    "alignment": "center",
                },
                "rw_2": {
                    "font_size": 9,
                    "font_name": "Arial",
                    "alignment": "center",
                },
                "rw_3": {
                    "font_size": 9,
                    "font_name": "Arial",
                    "font_color": (0, 118, 214),
                    "alignment": "center",
                }
            },
            "3monthsAsset": "5.6%",
            "1yrAssets": "7.7%",
            "assetsInception": "14.0%",
            "3monthsliab": "3.1%",
            "1yrliabs": "-6.7%",
            "liabsInception": "-3.6%",
            "3monthsOutPerformance": "2.4%",
            "1yrpr": "14.4%",
            "prInception": "17.6%",
        },
        "ACperformance":{
            "styles" : {
                "rw_1": {
                    "font_color": (252, 250, 250),
                    "bold": True,
                    "background_color": (14, 99, 179)
                },
                "rw_2": {
                    "bold": True,
                    "background_color": (46, 197, 217)
                },
                "rw_8": {
                    "bold": True,
                    "background_color": (46, 197, 217)
                },
                "rw_10": {
                    "font_color": (252, 250, 250),
                    "bold": True,
                    "background_color": (14, 99, 179)
                },
                "all": {
                    "font_size": 9,
                    "font_name": "Arial",
                    "alignment": "center",
                    "background_color": (252, 250, 250)
                },
            },
            "rows": [
                ["Scheme Performance", "5.6%" ,"7.7%"], 
                ["Total Growth","4.0%","15.3%"],
                ["Equities","5.4%","28.0%"],
                ["Corporate Bonds","2.6%","-0.1%"],
                ["Sovereign Bonds", "2.0%", "-1.4%"],
                ["Alternatives","4.8%","12.3%"],
                ["Dynamic Strategies","2.0%","13.9%"],
                ["Total Matching","13.3%","-24.7%"],
                ["LDI Funds & Cash","13.3%","-24.7%"],
                ["Liability Benchmark","3.1%","-6.7%"],
            ]
        },
        "assetAllocation":{
            "styles" : {
                "all":{
                    "font_size": 8.5,
                    "font_name": "Arial",
                    "font_color": (163, 162, 162),
                },
                "cl_0": {
                    "bold": True,
                    "alignment": "middle",
                },
                "rw_1": {
                    "column_indexes": [0, 1],
                    "font_color": (6, 123, 191),
                },
                "rw_2": {
                    "column_indexes": [0,1],
                    "font_color": (209, 189, 13),
                },
                "rw_3": {
                    "column_indexes": [0,1],
                    "font_color": (5, 125, 51),
                },
                "rw_4": {
                    "column_indexes": [0,1],
                    "font_color": (176, 19, 11),
                }
            },
            "data": {
                "eq": "EQUITIES",
                "eq_values": [
                    { "label": "UK Equity", "march": "1.8%", "june": "1.8%" },
                    { "label": "North America Equity", "march": "1.8%", "june": "1.9%" },
                    { "label": "Europe (ex UK) Equity", "march": "3.0%", "june": "2.7%" },
                    { "label": "Japan Equity", "march": "2.9%", "june": "2.8%" },
                    { "label": "Asia Pacific ex-Japan Equity", "march": "1.7%", "june": "1.7%" },
                    { "label": "Emerging Markets Equity", "march": "3.4%", "june": "3.6%" },
                    { "label": "Global Developed Small Cap Equity", "march": "2.7%", "june": "2.7%" },
                    { "label": "Smart beta equity", "march": "5.9%", "june": "5.6%" },
                ],
                "glb_bonds":"GLOBAL BONDS",
                "glb_values": [
                    { "label": "Fallen Angels’ Credit", "march": "5.1%", "june": "5.7%" },
                    { "label": "UK Investment Grade Credit", "march": "3.1%", "june": "4.9%" },
                    { "label": "Euro Investment Grade Credit", "march": "0.3%", "june": "0.7%" },
                    { "label": "US Investment Grade Credit", "march": "2.2%", "june": "2.2%" },
                    { "label": "Overseas Government Bonds", "march": "7.2%", "june": "5.5%" },
                ],
                "alternatives":"ALTERNATIVES",
                "alt_values": [
                    { "label": "Property", "march": "3.4%", "june": "3.4%" },
                    { "label": "Listed Private Equity", "march": "1.8%", "june": "1.7%" },
                    { "label": "High Yield Bonds", "march": "4.0%", "june": "3.9%" },
                    { "label": "Listed Infrastructure", "march": "4.1%", "june": "3.6%" },
                    { "label": "Global REITs", "march": "4.5%", "june": "4.4%" },
                    { "label": "Emerging Market Bonds (Local)", "march": "4.1%", "june": "4.2%" },
                    { "label": "Emerging Market Bonds (USD)", "march": "1.8%", "june": "1.9%" },
                    { "label": "Commodities", "march": "1.4%", "june": "1.5%" },
                ],
                
                "dynamic_str":"DYNAMIC  STRATEGIES",
                "dyn_str_values": [
                    { "label": "Multi-Asset Target Return (MATR)", "march": "16.5%", "june": "16.2%" },
                ],
                "lia_match":"LIABILITY-MATCHING",
                "lia_values":[
                    { "label": "Liability-matching credit", "march": "0.0%", "june": "0.0%" },
                    { "label": "Liability driven investment strategies", "march": "17.4%", "june": "17.5%" },
                ],
                "ratios":[
                    { "label": "Interest rate hedge ratio*", "march": "74%", "june": "73%" },
                    { "label": "Inflation hedge ratio*", "march": "74%", "june": "73%" },
                ]
            }
        },
        "toc":[
            {
                "id": "im1",
                "text": "This is a sample image",
                "sub":[{
                    "id": "im2",
                    "text": "sample image"
                    }]
                
            },
            {
                "id": "im3",
                "text": "Sample table to delete",
            }
        ],

        
        "position": "SSE",
        "city": "NW",
        "image_title": "This is a sample image",
        "sample_image": { "url" : "Sample-image.png" , "size": {"left":1,"top":1, "height":3, "width":8}},
        "project_description": "React , Node , AWS serverless",
        'table_name': "Sample table to delete",
        "remove_table_1": True,
        'table_name_row': "Sample table to delete row",
        "table_1_row_3_present": False,
        'table_name_column': "Sample table to delete column",
        "table_1_col_4_present": False,
        "sample_name": "Loop sample data",
        "sample_data_1": [{"name": "Kamal", "age": 12},{"name": "Amal", "age": 22},{"name": "Nuwan", "age": 32}],
        "sample_data_2": [{"name": "Sama", "age": 12},{"name": "Amara", "age": 22},{"name": "Nayana", "age": 32}],
        "sample_data_3": [{"city": "Colombo", "number": 1},{"city": "Colombo", "number": 2},{"city": "Colombo", "number": 3}],
        "cashFlows":{
                "headers": ["cashflow year","cashflow fixed","cashflow real"],
                "row_count": 10,
                "colum_count": 3,
                "table_count_per_slide": 4,
                "styles": {
                    "top": 1,
                    "width": 0.6,
                    "row_height": 0.04,
                },
                "rows": [ 
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"2","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"3","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"4","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"6","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"7","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"8","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"9","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"10","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"11","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"12","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"13","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"14","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"15","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"16","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"17","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"18","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"19","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"20","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"21","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"22","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"23","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"24","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"25","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"26","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"27","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"28","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"29","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"30","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"31","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"32","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"33","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"34","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"35","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"36","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"37","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"38","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"39","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"40","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"41","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"42","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"43","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"44","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"45","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"51","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"52","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"53","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"54","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"55","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"56","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"57","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"58","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"59","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"510","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"151","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"512","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"153","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"154","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"155","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"61","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"62","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"63","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"64","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"65","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"66","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"67","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"68","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"69","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"170","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"171","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"172","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"173","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"174","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"175","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"17","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"82","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"83","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"84","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"85","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"86","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"87","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"88","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"89","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"190","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"191","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"192","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"193","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"194","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"195","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"17","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"82","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"83","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"84","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"85","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"86","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"87","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"88","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"89","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"190","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"191","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"192","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"193","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"194","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"195","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"2","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"3","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"4","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"6","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"7","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"8","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"9","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"10","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"11","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"12","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"13","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"14","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"15","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"16","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"17","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"18","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"19","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"20","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"21","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"22","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"23","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"24","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"25","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"26","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"27","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"28","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"29","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"30","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"31","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"32","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"33","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"34","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"35","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"36","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"37","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"38","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"39","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"40","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"41","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"42","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"43","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"44","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"45","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"51","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5112","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5113","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5114","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5115","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5116","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5117","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"58","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5911","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"51110","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"15111","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"5112","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1513","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1514","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1515","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"611","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"612","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"6113","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"614","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"615","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"616","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"617","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"618","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"619","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1170","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1171","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1172","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1173","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1714","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1715","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1117","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"812","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"813","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"814","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"815","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"816","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"817","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"818","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"819","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1910","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1911","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1912","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1913","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1914","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1915","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"117","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"812","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"813","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"814","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"815","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"816","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"187","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"818","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"189","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1190","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1911","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1192","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1913","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1941","cashflow_real": "123123"
                    },
                    {
                        "cashflow_year": "2020","cashflow_fixed":"1915","cashflow_real": "123123"
                    }
                 ]
            }

    }

    slides = [slide for slide in prs.slides]
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # print("shape.text",shape.text)

                if("+++IF" in shape.text):
                    _if(prs, slide, shape,slides.index(slide), replacements)

                elif("+++FOR" in shape.text):
                    looper(prs, slide, shape,slides.index(slide), replacements)

                elif("+++IM" in shape.text):
                    replace_images(slide, shape, replacements)

                elif("+++TB_ADD" in shape.text):
                    replace_tables(prs, slide, shape,slides.index(slide), replacements)

                elif("+++TB_TX_UP" in shape.text):
                    update_table_text(prs, slide, shape,slides.index(slide), replacements)

                elif("+++TB_DRW" in shape.text):
                    drow_tables(prs, slide, shape,slides.index(slide), replacements)

                elif("+++INS" in shape.text):
                    text_replace(slide, shape, replacements)

    pres_final = remove_extra_slides(prs)
    drow_toc(prs,replacements)
                    

    prs.save('output.pptx')
    end_time = time.perf_counter ()
    print(end_time - start_time, "seconds")

    
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import math
import re
import pydash
from lxml import etree
from copy import deepcopy
from pptx.table import Table, _Row, _Column, _Cell
from pptx.enum.text import PP_ALIGN

from helpers.texts import text_tag_update
from helpers.utils import check_tag_exist, replace_tags, get_tag_content, get_tag_from_string, eval_executor, is_extra_slide

def replace_tables(presentation, slide, shape, slide_index, replacements):
    pattern = r'\+\+\+TB_ADD (.*?) \+\+\+'
    matches = get_tag_content(pattern, shape)

    if( not matches or len(matches) < 1):
        return

    for match in matches:
        object_value = pydash.get(replacements, match)
        if(object_value):
            replace_tags(str(f"+++TB_ADD {match} +++"), "", shape)
            create_table(presentation, slide, shape, slide_index, object_value)

def create_table(presentation, slide, shape, slide_index, replacements):
    row_count = replacements["row_count"]
    cols = replacements["colum_count"]
    headers = replacements["headers"]
    row_data = replacements["rows"]
    styles = replacements["styles"]
    table_count_per_slide = replacements["table_count_per_slide"]

    total_rows = len(row_data)
    total_table_count = math.ceil(total_rows / row_count)

    slide_count = math.ceil(total_table_count / table_count_per_slide) 
    extra_slide_count = slide_count - 1

    s = 0
    end = 0
    current_slide = 0
    while s < slide_count:
        j = 0
        left = 1

        if total_table_count > table_count_per_slide :
            table_count = table_count_per_slide
        else:
            table_count = total_table_count

        slide_row_start = 0
        slide_row_end = 0
        extra_slide_exists = False

        if(s > 0 and (not presentation.slides[slide_index + s])):
            break

        if(s > 0 and current_slide != s):
            for new_slide_shape in presentation.slides[slide_index + s].shapes:
                extra_slide_exists = is_extra_slide(presentation, slide_index + s, True)
                if(extra_slide_exists):
                    break
            current_slide = s

        if(s > 0 and (not extra_slide_exists)):
            return 
            

        while j < table_count:
            shape = presentation.slides[slide_index + s].shapes.add_table(row_count + 1, cols, Inches(left) , Inches(styles["top"]), Inches(styles["width"]), Inches(styles["row_height"]))
            table = shape.table
            tables_headers(table, headers)

            slide_row_start = (row_count * j ) + 1 + end
            slide_row_end = slide_row_start + (row_count - 1)

            tables_rows(table,row_data, slide_row_start, slide_row_end ,total_rows )
            left += 2
            j += 1

        end = slide_row_end
        total_table_count -= table_count_per_slide
        s += 1

def tables_headers(table, headers):
    i = 0
    for header in headers:
        table.cell(0, i).text = header
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(200, 253, 251)
        _set_cell_border(cell,"949595", '12000')
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(5)
        para.font.name = 'Comic Sans MS'
        para.font.color.rgb = RGBColor(19, 170, 246)
        table.columns[i].width = Inches(0.6)
        i += 1

def tables_rows(table, rowData, start, end,totalRows):
    j = start
    cell_start = 1

    entCount = end
    if end > totalRows:
        entCount = totalRows

    while j < entCount + 1:
        element = rowData[j - 1]
        k = 0
        if element:
            for key, value in element.items():
                table.cell(cell_start, k).text = value
                cell = table.cell(cell_start, k)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(228, 228, 228)
                _set_cell_border(cell,"949595", '12000')
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(6)
                para.font.name = 'Comic Sans MS'
                k += 1
        j += 1
        cell_start += 1

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

def remove_row(table, row):
    tbl = table._tbl
    tr = row._tr
    tbl.remove(tr)

def remove_tables(slide,content):
    #remove tables
    table_remove_pattern = r'\+\+\+TABLE_REMOVE (.*?) \+\+\+'
    table_remove_matches = get_tag_from_string(table_remove_pattern, content)
    
    if( table_remove_matches and len(table_remove_matches) > 0):
        table_remove_index_matches = table_remove_matches[0]
        table_id_tag = str(f"+++TB_ID {table_remove_index_matches} +++")
        _shap_count = 0
        for _shape in slide.shapes:
            if _shape.has_table: 
                for row in _shape.table.rows:
                    for cell in row.cells:
                        if table_id_tag in cell.text:
                            old_picture = slide.shapes[_shap_count]
                            old_pic = old_picture._element
                            old_pic.getparent().remove(old_pic)
                            break
            _shap_count += 1

def remove_table_rows(slide,content):
    table_row_remove_pattern = r'\+\+\+TABLE_ROW_REMOVE (.*?) \+\+\+'
    table_row_remove_matches = get_tag_from_string(table_row_remove_pattern, content)
    if( table_row_remove_matches and len(table_row_remove_matches) > 0):
        table_row_remove_index_matches = table_row_remove_matches[0]
        table_rw_id_tag = str(f"+++RW_ID {table_row_remove_index_matches} +++")
        for _shape in slide.shapes:
            if _shape.has_table: 
                for row_idx, row in enumerate(_shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if table_row_remove_index_matches in cell.text:
                            row_deleted = _shape.table.rows[row_idx]
                            remove_row(_shape.table, row_deleted)
                            break

def remove_table_column(slide,content):
    table_column_remove_pattern = r'\+\+\+TABLE_COLUMN_REMOVE (.*?) \+\+\+'
    table_column_remove_matches = get_tag_from_string(table_column_remove_pattern, content)
    if( table_column_remove_matches and len(table_column_remove_matches) > 0):
        table_column_remove_index_matches = table_column_remove_matches[0]
        for _shape in slide.shapes:
            if _shape.has_table:
                colum_index = ""
                for row_idx, row in enumerate(_shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if table_column_remove_index_matches in cell.text:
                            colum_index = col_idx
                            break

                for row_idx, row in enumerate(_shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if col_idx == colum_index:
                            cell._tc.delete()

                tree = etree.ElementTree(_shape.table._tbl)
                for e in tree.iter():
                    if(tree.getpath(e) == tree.getpath(_shape.table.columns[colum_index]._gridCol)):
                        e.getparent().remove(e)
                        break

def update_table_text(presentation, slide, shape, slide_index, replacements):
    pattern = r'\+\+\+TB_TX_UP (.*?) \+\+\+'
    matches = get_tag_content(pattern, shape)
    if( not matches or len(matches) < 1):
        return

    for match in matches:
        styles = False
        data = False
        data_path = False
        table_id = match

        if "DATA" in match:
            data_path = match.split(" DATA ")[1]
            table_id = match.split(" DATA ")[0]

        if data_path and data_path in replacements:
            data = replacements[data_path]

        if data and "styles" in data:
            styles = data["styles"]

        table_id_tag = str(f"+++TB_ID {table_id} +++")
        for _shape in slide.shapes:
            if _shape.has_table: 
                for row in _shape.table.rows:
                    for cell in row.cells:
                        if table_id_tag in cell.text:
                            execute_table_tags(_shape, _shape.table, data, styles)
                            new_text = cell.text.replace(str(f"+++TB_ID {table_id} +++"), "")
                            cell.text = new_text
                            break

        replace_tags(str(f"+++TB_TX_UP {match} +++"), "", shape)


def execute_table_tags(shape , table, data, styles):
    row_index = 0
    for row in table.rows:
        col_index = 0
        for cell in row.cells:
            pattern_for = r'\+\+\+FOR (.*?) FOR-END\+\+\+'
            matches_for = get_tag_from_string(pattern_for, cell.text)
            if( matches_for and len(matches_for) > 0):
                for match in matches_for:
                    pattern_condition = r'\(\((.*?)\)\)'
                    matched_condition = get_tag_from_string(pattern_condition,match)

                    pattern_content = r'\<\<(.*?)\>\>'
                    matched_content = get_tag_from_string(pattern_content,match)
                    for contidion in matched_condition:
                        object_value = pydash.get(data, contidion)
                        text_result = ""
                        if(object_value):
                            for _data in object_value:
                                updated_data = text_tag_update(matched_content[0],_data)
                                if(updated_data and updated_data["text"]):
                                    text_result += updated_data["text"] + "\n"
                        new_text = cell.text.replace(str(f"+++FOR {match} FOR-END+++"), text_result)
                        cell.text = new_text
                        try:
                            table_styles(cell,row_index,col_index,styles)
                        except ValueError:
                            print("error")
                

            pattern_text = r'\+\+\+INS (.*?) \+\+\+'
            matches_text_update = get_tag_from_string(pattern_text, cell.text)
            if( matches_text_update and len(matches_text_update) > 0):
                for match in matches_text_update:
                    new_text = cell.text.replace(str(f"+++INS {match} +++"), pydash.get(data, match))
                    cell.text = new_text
                    try:
                        table_styles(cell,row_index,col_index,styles)
                    except ValueError:
                        print("error")
            col_index += 1 
        row_index +=1     


def drow_tables(presentation, slide, shape, slide_index, replacements):
    pattern = r'\+\+\+TB_DRW (.*?) \+\+\+'
    matches = get_tag_content(pattern, shape)
    if( not matches or len(matches) < 1):
        return

    for match in matches:
        styles = False
        data = False
        data_path = False
        table_id = match

        if "DATA" in match:
            data_path = match.split(" DATA ")[1]
            table_id = match.split(" DATA ")[0]

        if data_path and data_path in replacements:
            data = replacements[data_path]

        if data and "styles" in data:
            styles = data["styles"]

        table_id_tag = str(f"+++TB_ID {table_id} +++")
        for _shape in slide.shapes:
            if _shape.has_table: 
                for row in _shape.table.rows:
                    for cell in row.cells:
                        if table_id_tag in cell.text:
                            execute_table_drower(_shape.table, data, styles)
                            new_text = cell.text.replace(str(f"+++TB_ID {table_id} +++"), "")
                            cell.text = new_text
                            break

        replace_tags(str(f"+++TB_DRW {match} +++"), "", shape)


def execute_table_drower(table, data,styles):
    row_data = data["rows"]
    row_index = 1
    for row in row_data:
        colum_index = 0
        if row_index > 1:
            add_new_row_to_existing_table(table)
        for column in row:
            cell = table.cell(row_index, colum_index)
            cell.text = column
            
            try:
                table_styles(cell,row_index,colum_index,styles)
            except ValueError:
                print("error")

            colum_index += 1
        row_index += 1

def add_new_row_to_existing_table(table):
    new_row = deepcopy(table._tbl.tr_lst[1])
    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = '' # defaulting cell contents to empty text

        table._tbl.append(new_row) 
        return table.rows[1]

def table_styles(cell,row_index,col_index,styles  ):
    try:
        row_st_index = str(f'rw_{row_index}')
        col_st_index = str(f'cl_{col_index}')
        para_index = 0
        for paragraph in cell.text_frame.paragraphs:
            para = cell.text_frame.paragraphs[para_index]
            if(styles and 'all' in styles):
                common_styles = styles['all']

                if('font_size' in common_styles):
                    para.font.size = Pt(common_styles['font_size'])
                if('font_name' in common_styles):
                    para.font.name = common_styles['font_name']
                if('bold' in common_styles):
                    para.font.bold = common_styles['bold']
                if('italic' in common_styles):
                    para.font.italic = common_styles['italic']
                if("font_color" in common_styles):
                    para.font.color.rgb = RGBColor(common_styles["font_color"][0], common_styles["font_color"][1],common_styles["font_color"][2])
                if("alignment" in common_styles):
                    if common_styles["alignment"] == "center":
                        para.alignment = PP_ALIGN.CENTER
                if("background_color" in common_styles):
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(common_styles["background_color"][0], common_styles["background_color"][1],common_styles["background_color"][2])

            if(styles and row_st_index in styles):
                _styles = styles[row_st_index]
                if( "column_indexes" in _styles):
                    if(col_index in _styles["column_indexes"]):
                        if('font_size' in _styles):
                            para.font.size = Pt(_styles['font_size'])
                        if('font_name' in _styles):
                            para.font.name = _styles['font_name']
                        if('bold' in _styles):
                            para.font.bold = _styles['bold']
                        if('italic' in _styles):
                            para.font.italic = _styles['italic']
                        if("font_color" in _styles):
                            para.font.color.rgb = RGBColor(_styles["font_color"][0], _styles["font_color"][1],_styles["font_color"][2])
                        if("background_color" in _styles):
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(_styles["background_color"][0], _styles["background_color"][1],_styles["background_color"][2])
                else:
                    if('font_size' in _styles):
                        para.font.size = Pt(_styles['font_size'])
                    if('font_name' in _styles):
                        para.font.name = _styles['font_name']
                    if('bold' in _styles):
                        para.font.bold = _styles['bold']
                    if('italic' in _styles):
                        para.font.italic = _styles['italic']
                    if("font_color" in _styles):
                        para.font.color.rgb = RGBColor(_styles["font_color"][0], _styles["font_color"][1],_styles["font_color"][2])
                    if("background_color" in _styles):
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(_styles["background_color"][0], _styles["background_color"][1],_styles["background_color"][2])

            if(styles and col_st_index in styles):
                col_styles = styles[col_st_index]

                if('font_size' in col_styles):
                    para.font.size = Pt(col_styles['font_size'])
                if('font_name' in col_styles):
                    para.font.name = col_styles['font_name']
                if('bold' in col_styles):
                    para.font.bold = col_styles['bold']
                if('italic' in col_styles):
                    para.font.italic = col_styles['italic']
                if("font_color" in col_styles):
                    para.font.color.rgb = RGBColor(col_styles["font_color"][0], col_styles["font_color"][1],col_styles["font_color"][2])
                if("background_color" in col_styles):
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(col_styles["background_color"][0], col_styles["background_color"][1],col_styles["background_color"][2])
    
            para_index += 1
    except ValueError:
        print("error")
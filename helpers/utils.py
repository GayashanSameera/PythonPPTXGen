from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import math
import re
import pydash
from pptx.table import Table, _Row, _Column, _Cell
from copy import deepcopy

def check_tag_exist(tag, shape):
    matches = tag in shape.text
    return matches

def replace_tags(replaced_for,replaced_text, shape):
    if shape.has_text_frame:
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                cur_text = run.text
                print("cur_text",cur_text)
                new_text = cur_text.replace(replaced_for, replaced_text)
                print("new_text",new_text)
                run.text = new_text

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if replaced_for in cell.text:
                    new_text = cell.text.replace(replaced_for, replaced_text)
                    cell.text = new_text

def get_tag_content(pattern,shape):
    matches = re.findall(pattern, shape.text)
    return matches

def get_tag_from_string(pattern,string):
    matches = re.findall(pattern, string)
    return matches

def eval_executor(logic, replacements):
    return eval(logic,replacements)

def delete_slides(presentation, index):
    print("index",index)
    xml_slides = presentation.slides._sldIdLst  
    slides = list(xml_slides)
    try:
        slides[index]
        xml_slides.remove(slides[index])
    except ValueError:
        print("error") 

def is_extra_slide(presentation, slide_index, remove_tag):
    extra_slide_exists = False
    extra = 'EXTRA_SLIDE'
    for new_slide_shape in presentation.slides[slide_index].shapes:
        if new_slide_shape.has_text_frame:
            matches = check_tag_exist(extra, new_slide_shape)
            if(matches):
                extra_slide_exists = True
                if(remove_tag):
                    replace_tags(str(f"+++ {extra} +++"), "", new_slide_shape)
                break
    return extra_slide_exists

def remove_extra_slides(presentation):
    extra = 'EXTRA_SLIDE'
    slides = [slide for slide in presentation.slides]
    slide_indexs_to_delete = []
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                matches = check_tag_exist(extra, shape)
                if(matches):
                    slide_indexs_to_delete.append(slides.index(slide))

    if(len(slide_indexs_to_delete) > 0):
        array_index = 0
        for s_index in slide_indexs_to_delete:
            slide_index = s_index
            delete_slides(presentation, slide_index - array_index)
            array_index += 1 


    return presentation

def drow_toc(presentation,data):
    ids = identify_tags_with_page_number(presentation)
    print("ids",ids)
    slides = [slide for slide in presentation.slides]
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if("+++TOC" in shape.text):
                    pattern = r'\+\+\+TOC (.*?) \+\+\+'
                    matches = get_tag_content(pattern, shape)
                    print(matches)
                    data_path = matches[0]
                    datam = pydash.get(data, data_path)
                    print("datam",datam)
                    update_toc_table(slide,datam,ids)
                    replace_tags(str(f"+++TOC {data_path} +++"), "", shape)
                    return

def update_toc_table(slide,datam,ids):
    for shape in slide.shapes:
        if shape.has_table:
            execute_table_drower(shape.table, datam, ids)

def execute_table_drower(table, data,ids):
    row_index = 0
    for row in data:
        print("row",row)
        if row_index > 0:
            add_new_row_to_existing_table(table)
        
        cell_1 = table.cell(row_index, 0)
        cell_1.text = row["text"]
        cell_2 = table.cell(row_index, 2)
        row_id = row["id"]
        print("row_id",row_id)
        print("ids",ids)
        print("ids[row_id]",ids[row_id])
        cell_2.text = str(ids[row_id])

        row_index += 1

def add_new_row_to_existing_table(table):
    new_row = deepcopy(table._tbl.tr_lst[0])
    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = ''

        table._tbl.append(new_row) 
        return table.rows[0]
 

def identify_tags_with_page_number(presentation):
    res = {}
    slides = [slide for slide in presentation.slides]
    pattern = r'\+\+\+TOC_IDS (.*?) \+\+\+'
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                matches = get_tag_content(pattern, shape)
                if(matches and len(matches) > 0):
                    match_string = matches[0]
                    page_number = slides.index(slide) + 1
                    print("match_string",match_string)
                    print("page_number",page_number)
                    if "," in match_string:
                        id_array = match_string.split(",")
                        print("id_array",id_array)
                        for id in id_array:
                            res[id] = page_number
                    else:
                        res[match_string] = page_number

            replace_tags(str(f"+++TOC_IDS {matches} +++"), "", shape)
    print("res",res)
    return res

